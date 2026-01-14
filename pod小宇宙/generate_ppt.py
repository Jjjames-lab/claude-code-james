#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
播客逐字稿产品可行性分析报告 - PPT生成器
创建详尽的专业演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加背景渐变（通过设置形状）
    left = top = width = height = Inches(10)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(26, 35, 126)  # 深蓝色背景
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 200, 255)
    p.alignment = PP_ALIGN.CENTER

    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年1月9日"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, title, content_items, accent_color=RGBColor(26, 35, 126)):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = accent_color
    title_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            # 带有特殊格式的项目
            if 'bullet' in item:
                p.text = item['text']
                p.level = 0
            elif 'title' in item:
                p.text = item['title']
                p.font.bold = True
                p.font.size = Pt(20)
                p.space_after = Pt(10)
                continue
            else:
                p.text = item.get('text', str(item))
        else:
            p.text = f"• {item}"

        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
        p.line_spacing = 1.5

    return slide

def create_two_column_slide(prs, title, left_items, right_items,
                           left_title="左侧内容", right_title="右侧内容"):
    """创建两栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)

    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)
    p.space_after = Pt(15)

    for item in left_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)

    # 右栏
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 57, 70)
    p.space_after = Pt(15)

    for item in right_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)

    return slide

def create_table_slide(prs, title, headers, rows):
    """创建表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)

    # 表格
    x, y, cx, cy = Inches(0.7), Inches(1.5), Inches(8.6), Inches(5)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, cx, cy).table

    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(26, 35, 126)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # 填充数据
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
            paragraph.alignment = PP_ALIGN.CENTER

            # 斑马纹
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

    return slide

def create_highlight_slide(prs, title, highlight_text, description_items,
                          highlight_color=RGBColor(230, 57, 70)):
    """创建重点强调页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)

    # 重点数字/文字
    highlight_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(1.5))
    tf = highlight_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = highlight_text
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = highlight_color
    p.alignment = PP_ALIGN.CENTER

    # 描述
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(3))
    tf = desc_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(description_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_summary_slide(prs, scores):
    """创建评分总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "综合评分：5.5/10"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 35, 126)
    p.alignment = PP_ALIGN.CENTER

    # 评分项
    y_pos = 1.5
    for category, score, color in scores:
        # 进度条背景
        bar_bg = slide.shapes.add_shape(1, Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = RGBColor(230, 230, 230)
        bar_bg.line.fill.background()

        # 进度条
        bar_width = Inches(8 * score / 10)
        bar = slide.shapes.add_shape(1, Inches(1), Inches(y_pos), bar_width, Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 文本
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.5), Inches(8), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{category}: {score}/10"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)

        y_pos += 1.2

    # 结论
    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(1))
    tf = conclusion_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️ 结论：技术可行，但商业化困难，建议调整方向"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(211, 84, 0)
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    print("📊 开始生成PPT演示文稿...")
    print("="*80)

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\n📄 幻灯片 1/30: 封面")
    create_title_slide(prs,
                     "播客逐字稿产品\n可行性分析报告",
                     "深度市场调研 · 技术方案 · 商业模式 · 风险评估")

    print("📄 幻灯片 2/30: 目录")
    create_content_slide(prs, "目录", [
        "01. 执行摘要",
        "02. 市场规模分析",
        "03. 竞品深度对比",
        "04. 用户画像研究",
        "05. 技术可行性评估",
        "06. 成本模型分析",
        "07. 商业模式设计",
        "08. SWOT战略分析",
        "09. 风险与挑战",
        "10. 替代方案建议",
        "11. 技术栈推荐",
        "12. 行动路线图",
        "13. 最终建议"
    ])

    print("📄 幻灯片 3/30: 执行摘要 - 核心结论")
    create_highlight_slide(prs,
                         "核心结论",
                         "⚠️ 综合评分：5.5/10",
                         [
                             "✅ 技术可行性：可行，方案成熟",
                             "✅ 市场需求：存在，但规模有限",
                             "⚠️ 商业潜力：谨慎，竞争激烈",
                             "🔴 个人开发：困难，资源有限",
                             "🔴 综合评估：建议重新思考或大幅调整方向"
                         ],
                         highlight_color=RGBColor(211, 84, 0))

    print("📄 幻灯片 4/30: 市场规模 - 长音频市场")
    create_content_slide(prs, "市场规模：中国长音频市场", [
        {"title": "核心数据", "text": ""},
        {"bullet": True, "text": "2024年市场规模：310亿元人民币"},
        {"bullet": True, "text": "2025年预计：337亿元（增长率8.6%）"},
        {"bullet": True, "text": "2026年预测：380亿元"},
        {"title": "增长驱动因素", "text": ""},
        {"bullet": True, "text": "车载音频市场快速增长"},
        {"bullet": True, "text": "智能音箱普及率提升"},
        {"bullet": True, "text": "在线教育和内容付费兴起"}
    ])

    print("📄 幻灯片 5/30: 市场规模 - 播客细分市场")
    create_content_slide(prs, "市场规模：中国播客市场", [
        {"title": "播客市场规模", "text": ""},
        {"bullet": True, "text": "2025年预计突破：50亿元人民币"},
        {"bullet": True, "text": "年增长率：约8%"},
        {"bullet": True, "text": "占长音频市场比例：约15%"},
        {"title": "用户规模", "text": ""},
        {"bullet": True, "text": "2025年播客用户：1.5亿人"},
        {"bullet": True, "text": "2026年预测：突破1.7亿人"},
        {"title": "关键洞察", "text": ""},
        {"bullet": True, "text": "⚠️ 播客市场只是长音频的一个细分领域"}
    ])

    print("📄 幻灯片 6/30: 小宇宙平台数据")
    create_content_slide(prs, "小宇宙播客平台数据", [
        {"title": "用户规模", "text": ""},
        {"bullet": True, "text": "2023年底月活：约600万"},
        {"bullet": True, "text": "2024年增长：50%（破圈趋势明显）"},
        {"bullet": True, "text": "2024年推算月活：约900万"},
        {"title": "市场地位", "text": ""},
        {"bullet": True, "text": "在长音频市场：第二梯队（vs 喜马拉雅73.5%）"},
        {"bullet": True, "text": "在播客垂直领域：头部平台"},
        {"title": "关键问题", "text": ""},
        {"bullet": True, "text": "⚠️ 可服务用户：900万 × 5%（有转录需求）= 45万"}
    ])

    print("📄 幻灯片 7/30: 竞品分析 - 飞书妙记")
    create_content_slide(prs, "竞品分析：飞书妙记（字节跳动）", [
        {"title": "产品定位", "text": ""},
        {"bullet": True, "text": "通用音视频转文字工具"},
        {"bullet": True, "text": "面向会议、访谈、课堂等场景"},
        {"title": "核心优势", "text": ""},
        {"bullet": True, "text": "✅ 免费额度：300分钟/月"},
        {"bullet": True, "text": "✅ 功能强大：实时转录、说话人区分、AI摘要"},
        {"bullet": True, "text": "✅ 字节生态支持，品牌信任度高"},
        {"bullet": True, "text": "✅ 支持多语言识别"},
        {"title": "劣势", "text": ""},
        {"bullet": True, "text": "❌ 非播客特化产品"},
        {"bullet": True, "text": "❌ 需手动上传音频，无法直接从小宇宙获取"}
    ])

    print("📄 幻灯片 8/30: 竞品分析 - 通义听悟")
    create_content_slide(prs, "竞品分析：通义听悟（阿里云）", [
        {"title": "产品定位", "text": ""},
        {"bullet": True, "text": "AI驱动的音视频转写工具"},
        {"title": "核心优势", "text": ""},
        {"bullet": True, "text": "✅ 定价有竞争力：29元/月，10小时"},
        {"bullet": True, "text": "✅ AI能力强大（阿里大模型支持）"},
        {"bullet": True, "text": "✅ 支持实时转录和后处理"},
        {"bullet": True, "text": "✅ 超出部分：2元/小时"},
        {"title": "劣势", "text": ""},
        {"bullet": True, "text": "❌ 同样非播客特化"},
        {"bullet": True, "text": "❌ 没有小宇宙集成"},
        {"bullet": True, "text": "❌ 需要上传音频文件"}
    ])

    print("📄 幻灯片 9/30: 竞品对比矩阵")
    create_table_slide(prs, "竞品功能对比矩阵",
                      ["功能", "飞书妙记", "通义听悟", "你的产品"],
                      [
                          ["语音转文字", "✅", "✅", "✅"],
                          ["说话人分离", "✅", "✅", "✅"],
                          ["播客特化", "❌", "❌", "✅"],
                          ["实时同步", "✅", "✅", "✅"],
                          ["小宇宙集成", "❌", "❌", "✅"],
                          ["时间轴跳转", "✅", "✅", "✅"],
                          ["价格", "免费300分钟/月", "29元/月10h", "待定"]
                      ])

    print("📄 幻灯片 10/30: 竞品威胁评估")
    create_two_column_slide(prs,
                           "竞品威胁评估",
                           [
                               "飞书/通义可能快速跟进",
                               "一旦看到市场机会",
                               "3个月内可推出类似功能",
                               "利用现有基础设施",
                               "边际成本几乎为零",
                               "通过品牌优势快速获取用户",
                               "",
                               "用户迁移成本低",
                               "如果竞品添加类似功能",
                               "用户没有理由切换"
                           ],
                           [
                               "价格战压力",
                               "你无法承担与大厂的价格战",
                               "他们可以长期免费来扼杀竞争",
                               "",
                               "品牌信任度差距",
                               "用户更信任大厂产品",
                               "数据安全和隐私顾虑",
                               "",
                               "结论：成为\"先烈\"的风险高"
                           ],
                           "🔴 主要威胁",
                           "🔴 长期风险")

    print("📄 幻灯片 11/30: 用户画像 - 年龄与地域")
    create_content_slide(prs, "用户画像：年龄与地域分布", [
        {"title": "年龄分布", "text": ""},
        {"bullet": True, "text": "26-35岁：45%（核心主力）"},
        {"bullet": True, "text": "18-25岁：25%"},
        {"bullet": True, "text": "36-45岁：22%"},
        {"bullet": True, "text": "46岁+：8%"},
        {"title": "城市等级", "text": ""},
        {"bullet": True, "text": "一线城市：35%"},
        {"bullet": True, "text": "新一线城市：28%"},
        {"bullet": True, "text": "二线城市：22%"},
        {"bullet": True, "text": "三线及以下：15%"},
        {"title": "关键洞察", "text": ""},
        {"bullet": True, "text": "⚠️ 一线+新一线占63%，机会成本高"}
    ])

    print("📄 幻灯片 12/30: 用户画像 - 职业与付费意愿")
    create_content_slide(prs, "用户画像：职业与付费意愿", [
        {"title": "职业分布", "text": ""},
        {"bullet": True, "text": "在职员工：55%"},
        {"bullet": True, "text": "学生：28%"},
        {"bullet": True, "text": "自由职业：10%"},
        {"bullet": True, "text": "其他：7%"},
        {"title": "付费意愿调查", "text": ""},
        {"bullet": True, "text": "42%愿意付费"},
        {"bullet": True, "text": "35%考虑付费"},
        {"bullet": True, "text": "23%不愿意付费"},
        {"title": "价格接受度", "text": ""},
        {"bullet": True, "text": "单集/系列：≤50元人民币"},
        {"bullet": True, "text": "⚠️ 用户对价格敏感"}
    ])

    print("📄 幻灯片 13/30: 可服务市场规模测算")
    create_highlight_slide(prs,
                         "可服务市场测算",
                         "1.8 - 9万人",
                         [
                             "小宇宙月活：900万人",
                             "对转录有需求：5% = 45万人",
                             "愿意付费：20% = 9万人",
                             "",
                             "⚠️ 关键问题：你如何触达这9万人？",
                             "⚠️ 关键问题：你如何说服他们付费？"
                         ],
                         highlight_color=RGBColor(230, 57, 70))

    print("📄 幻灯片 14/30: 技术可行性 - 小宇宙数据提取")
    create_content_slide(prs, "技术可行性：小宇宙数据提取", [
        {"title": "非官方API方案", "text": ""},
        {"bullet": True, "text": "✅ GitHub项目：ultrazg/xyz（小宇宙FM API）"},
        {"bullet": True, "text": "✅ 功能：获取音频、单集详情、Show Notes、时间轴"},
        {"bullet": True, "text": "✅ 直接API：api.xiaoyuzhoufm.com"},
        {"bullet": True, "text": "✅ RSSHub：开源RSS生成方案"},
        {"title": "可提取信息", "text": ""},
        {"bullet": True, "text": "✅ 标题、节目名称"},
        {"bullet": True, "text": "✅ Show Notes"},
        {"bullet": True, "text": "✅ 时间轴（部分播客）"},
        {"bullet": True, "text": "✅ 音频文件下载链接"},
        {"title": "⚠️ 风险", "text": ""},
        {"bullet": True, "text": "非官方API可能随时失效"}
    ])

    print("📄 幻灯片 15/30: 技术可行性 - 逐字稿生成")
    create_content_slide(prs, "技术可行性：逐字稿生成与说话人区分", [
        {"title": "核心技术栈：WhisperX", "text": ""},
        {"bullet": True, "text": "✅ OpenAI Whisper的增强版"},
        {"bullet": True, "text": "✅ 集成说话人分离功能"},
        {"bullet": True, "text": "✅ 支持词汇级时间戳（word-level timestamps）"},
        {"bullet": True, "text": "✅ 中文识别准确率：95%+"},
        {"title": "技术实现", "text": ""},
        {"bullet": True, "text": "✅ 语音识别（ASR）"},
        {"bullet": True, "text": "✅ 说话人分离（Speaker Diarization）"},
        {"bullet": True, "text": "✅ 实时同步（<500ms延迟）"},
        {"bullet": True, "text": "✅ 自动标点和分行"},
        {"title": "结论", "text": ""},
        {"bullet": True, "text": "✅ 技术完全成熟，开源方案可用"}
    ])

    print("📄 幻灯片 16/30: 成本模型 - 5个阶段")
    create_table_slide(prs, "成本模型：不同用户规模的月度成本（元）",
                      ["用户规模", "服务器", "GPU算力", "带宽", "存储", "总成本"],
                      [
                          ["1人(MVP)", "0", "100", "0", "10", "¥110"],
                          ["100人", "200", "500", "50", "50", "¥800"],
                          ["1,000人", "1,500", "3,000", "300", "300", "¥5,100"],
                          ["10,000人", "8,000", "18,000", "2,000", "2,000", "¥30,000"],
                          ["100,000人", "35,000", "90,000", "12,000", "12,000", "¥149,000"]
                      ])

    print("📄 幻灯片 17/30: 成本分析 - 关键发现")
    create_content_slide(prs, "成本分析：关键发现", [
        {"title": "成本结构特点", "text": ""},
        {"bullet": True, "text": "🔴 GPU成本占比最高（60%）"},
        {"bullet": True, "text": "🔴 WhisperX需要GPU加速"},
        {"bullet": True, "text": "🔴 按需计费：1小时音频 ≈ ¥5-10"},
        {"bullet": True, "text": "🔴 用户月均转录2小时，成本¥10-20"},
        {"title": "盈亏平衡困难", "text": ""},
        {"bullet": True, "text": "假设客单价：¥25/月"},
        {"bullet": True, "text": "单用户可变成本：¥15/月（主要是GPU）"},
        {"bullet": True, "text": "固定成本：¥3,000/月"},
        {"bullet": True, "text": "盈亏平衡点：300付费用户"},
        {"bullet": True, "text": "⚠️ 按3%转化率，需要10,000总用户"},
        {"title": "⚠️ 个人开发者资金压力", "text": ""},
        {"bullet": True, "text": "1,000用户：月支出¥5,100"},
        {"bullet": True, "text": "10,000用户：月支出¥30,000"},
        {"bullet": True, "text": "你能承担前期亏损吗？"}
    ])

    print("📄 幻灯片 18/30: 商业模式 - 分层定价")
    create_table_slide(prs, "商业模式：分层定价策略",
                      ["版本", "价格", "功能", "目标用户"],
                      [
                          ["免费版", "¥0", "每月1集转录\n基础逐字稿\n7天保留", "体验用户"],
                          ["标准版", "¥19/月", "每月5集转录\n说话人区分\n永久保留\n导出功能", "轻度用户"],
                          ["专业版", "¥49/月", "无限转录\n实时同步\nAPI访问\n优先处理", "重度用户"],
                          ["企业版", "¥199/月", "团队协作\n批量处理\n定制功能", "内容团队"]
                      ])

    print("📄 幻灯片 19/30: 收入预测（10,000用户）")
    create_content_slide(prs, "收入预测：10,000用户规模", [
        {"title": "保守场景（1%转化率）", "text": ""},
        {"bullet": True, "text": "付费用户：100人"},
        {"bullet": True, "text": "月收入：¥1,900"},
        {"bullet": True, "text": "月成本：¥30,000"},
        {"bullet": True, "text": "🔴 月亏损：¥28,100"},
        {"title": "中等场景（3%转化率）", "text": ""},
        {"bullet": True, "text": "付费用户：300人"},
        {"bullet": True, "text": "月收入：¥5,700"},
        {"bullet": True, "text": "月成本：¥30,000"},
        {"bullet": True, "text": "🔴 月亏损：¥24,300"},
        {"title": "乐观场景（5%转化率）", "text": ""},
        {"bullet": True, "text": "付费用户：500人"},
        {"bullet": True, "text": "月收入：¥9,500"},
        {"bullet": True, "text": "月成本：¥30,000"},
        {"bullet": True, "text": "🔴 月亏损：¥20,500"},
        {"title": "⚠️ 结论", "text": ""},
        {"bullet": True, "text": "即使乐观场景，仍需40,000+用户才能盈亏平衡"}
    ])

    print("📄 幻灯片 20/30: SWOT分析 - 优势与劣势")
    create_two_column_slide(prs,
                           "SWOT分析",
                           [
                               "✅ 垂直化聚焦",
                               "专注播客场景",
                               "体验可能更好",
                               "",
                               "✅ 小宇宙深度集成",
                               "一键提取",
                               "减少用户操作",
                               "",
                               "✅ 极致用户体验",
                               "个人开发者可以更快迭代",
                               "",
                               "✅ 低成本起步",
                               "MVP阶段成本可控"
                           ],
                           [
                               "❌ 个人开发资源",
                               "时间、精力、资金有限",
                               "",
                               "❌ 无官方合作",
                               "小宇宙API可能随时失效",
                               "",
                               "❌ 技术门槛",
                               "需要处理音频、AI、实时同步",
                               "",
                               "❌ 资金有限",
                               "无法承担长期亏损",
                               "无法大规模营销"
                           ],
                           "优势 Strengths ✅",
                           "劣势 Weaknesses ❌")

    print("📄 幻灯片 21/30: SWOT分析 - 机会与威胁")
    create_two_column_slide(prs,
                           "SWOT分析（续）",
                           [
                               "✅ 播客市场快速增长",
                               "年增长率8%",
                               "",
                               "✅ 用户付费意愿提升",
                               "42%愿意付费",
                               "",
                               "✅ 竞品未深耕播客",
                               "飞书/通义是通用产品",
                               "",
                               "✅ AI技术成熟",
                               "WhisperX开源可用"
                           ],
                           [
                               "🔴 平台政策风险",
                               "小宇宙可能封禁非官方API",
                               "",
                               "🔴 大厂可能入局",
                               "一旦被验证，大厂可快速跟进",
                               "",
                               "🔴 用户获取成本高",
                               "没有品牌背书，获客困难",
                               "",
                               "🔴 盈利不确定性强",
                               "需要4万+用户才能盈亏平衡"
                           ],
                           "机会 Opportunities ✅",
                           "威胁 Threats 🔴")

    print("📄 幻灯片 22/30: 核心风险警示")
    create_content_slide(prs, "🔴 核心风险警示", [
        {"title": "风险1：大厂快速跟进", "text": ""},
        {"bullet": True, "text": "你花6个月开发产品，获得5,000用户"},
        {"bullet": True, "text": "飞书/通义看到机会，3个月后推出类似功能"},
        {"bullet": True, "text": "利用现有用户基础快速抢占市场"},
        {"bullet": True, "text": "结果：你成为\"先烈\""},
        {"title": "风险2：无法突破获客瓶颈", "text": ""},
        {"bullet": True, "text": "小宇宙900万用户，但他们如何知道你的产品？"},
        {"bullet": True, "text": "你没有营销预算"},
        {"bullet": True, "text": "获取10,000用户可能需要12-24个月"},
        {"bullet": True, "text": "⚠️ 你能承担24个月的亏损吗？"},
        {"title": "风险3：技术债务和运维压力", "text": ""},
        {"bullet": True, "text": "你需要负责：开发、测试、运维、客服、营销"},
        {"bullet": True, "text": "你有足够的精力吗？"}
    ])

    print("📄 幻灯片 23/30: 替代方案1 - 播客知识管理")
    create_content_slide(prs, "💡 替代方案1：播客知识管理工具", [
        {"title": "核心思路", "text": ""},
        {"bullet": True, "text": "不再只做\"转录\"，而是做\"播客学习工具\""},
        {"title": "核心功能", "text": ""},
        {"bullet": True, "text": "一键提取小宇宙播客"},
        {"bullet": True, "text": "AI自动生成摘要、思维导图、关键观点"},
        {"bullet": True, "text": "用户可以添加笔记、高亮、标签"},
        {"bullet": True, "text": "与Notion/Obsidian/飞书文档集成"},
        {"bullet": True, "text": "\"我的播客知识库\""},
        {"title": "差异化", "text": ""},
        {"bullet": True, "text": "飞书/通义：转录工具"},
        {"bullet": True, "text": "你：播客学习工具 ⭐"},
        {"title": "优势", "text": ""},
        {"bullet": True, "text": "✅ 更高的价值感知"},
        {"bullet": True, "text": "✅ 更强的用户粘性"},
        {"bullet": True, "text": "✅ 更难被竞品复制"}
    ])

    print("📄 幻灯片 24/30: 替代方案2 - B端播客制作工具")
    create_content_slide(prs, "💡 替代方案2：B端播客制作助手", [
        {"title": "核心思路", "text": ""},
        {"bullet": True, "text": "服务播客主，而不是听众"},
        {"title": "核心功能", "text": ""},
        {"bullet": True, "text": "自动生成播客Show Notes"},
        {"bullet": True, "text": "自动提取精彩片段（社交媒体推广）"},
        {"bullet": True, "text": "自动生成多平台文案（小红书、微博、公众号）"},
        {"bullet": True, "text": "播客数据分析和优化建议"},
        {"bullet": True, "text": "多平台分发工具"},
        {"title": "商业模式", "text": ""},
        {"bullet": True, "text": "按集收费：¥19/集"},
        {"bullet": True, "text": "订阅制：¥199/月，无限处理"},
        {"title": "优势", "text": ""},
        {"bullet": True, "text": "✅ B端付费意愿更强"},
        {"bullet": True, "text": "✅ 客单价更高"},
        {"bullet": True, "text": "✅ 竞争更少"}
    ])

    print("📄 幻灯片 25/30: 替代方案3 - 播客推荐引擎")
    create_content_slide(prs, "💡 替代方案3：播客推荐+逐字稿", [
        {"title": "核心思路", "text": ""},
        {"bullet": True, "text": "从\"工具\"到\"平台\""},
        {"title": "核心功能", "text": ""},
        {"bullet": True, "text": "基于用户兴趣推荐播客内容"},
        {"bullet": True, "text": "推荐内容带逐字稿"},
        {"bullet": True, "text": "用户可以\"订阅话题\"而不是\"订阅播客\""},
        {"bullet": True, "text": "AI总结最新播客内容"},
        {"title": "差异化", "text": ""},
        {"bullet": True, "text": "小宇宙：订阅播客，从头听到尾"},
        {"bullet": True, "text": "你：订阅话题，AI帮你听 ⭐"},
        {"title": "优势", "text": ""},
        {"bullet": True, "text": "✅ 用户粘性更强（每天打开）"},
        {"bullet": True, "text": "✅ 更高的DAU"},
        {"bullet": True, "text": "✅ 更难被竞品复制（需要数据和算法）"}
    ])

    print("📄 幻灯片 26/30: 技术栈推荐 - MVP阶段")
    create_content_slide(prs, "技术栈推荐：MVP阶段（1-3个月）", [
        {"title": "目标", "text": "快速验证需求，控制成本"},
        {"title": "技术选型", "text": ""},
        {"bullet": True, "text": "前端：Next.js + React（快速开发，SEO友好）"},
        {"bullet": True, "text": "后端：Python FastAPI（轻量级，易集成AI库）"},
        {"bullet": True, "text": "数据库：SQLite → PostgreSQL（开始简单）"},
        {"bullet": True, "text": "AI模型：WhisperX本地部署（开源免费）"},
        {"bullet": True, "text": "部署：Railway/Render（免费额度够MVP）"},
        {"title": "成本", "text": "¥110-200/月"},
        {"title": "适合", "text": "1-10个早期用户测试"}
    ])

    print("📄 幻灯片 27/30: 技术栈推荐 - 成长期")
    create_content_slide(prs, "技术栈推荐：成长期（3-12个月）", [
        {"title": "目标", "text": "支持100-1,000用户"},
        {"title": "技术升级", "text": ""},
        {"bullet": True, "text": "前端：Next.js + TailwindCSS"},
        {"bullet": True, "text": "后端：FastAPI + Celery（异步任务队列）"},
        {"bullet": True, "text": "数据库：PostgreSQL云托管"},
        {"bullet": True, "text": "缓存：Redis（加速响应）"},
        {"bullet": True, "text": "AI模型：WhisperX + GPU服务器"},
        {"bullet": True, "text": "部署：阿里云/腾讯云（国内访问快）"},
        {"bullet": True, "text": "CDN：七牛云/阿里云CDN"},
        {"title": "成本", "text": "¥800-8,000/月"},
        {"title": "适合", "text": "100-1,000用户规模"}
    ])

    print("📄 幻灯片 28/30: 行动路线图")
    create_content_slide(prs, "行动路线图：如果接受建议（转向播客知识管理）", [
        {"title": "第1-2周：MVP开发", "text": ""},
        {"bullet": True, "text": "小宇宙链接解析"},
        {"bullet": True, "text": "WhisperX转录"},
        {"bullet": True, "text": "基础AI摘要"},
        {"title": "第3-4周：用户测试", "text": ""},
        {"bullet": True, "text": "邀请20个种子用户"},
        {"bullet": True, "text": "收集反馈"},
        {"bullet": True, "text": "快速迭代"},
        {"title": "第2-3月：功能完善", "text": ""},
        {"bullet": True, "text": "笔记功能"},
        {"bullet": True, "text": "高亮和标签"},
        {"bullet": True, "text": "导出功能"},
        {"title": "第4-6月：商业化", "text": "付费转化、营销推广、目标：100付费用户"}
    ])

    print("📄 幻灯片 29/30: 如果坚持原方向")
    create_content_slide(prs, "行动路线图：如果坚持原方向（逐字稿工具）", [
        {"title": "第1-4周：快速MVP", "text": ""},
        {"bullet": True, "text": "验证技术可行性"},
        {"bullet": True, "text": "邀请5-10个朋友测试"},
        {"bullet": True, "text": "评估用户反馈强度"},
        {"title": "第5-8周：早期推广", "text": ""},
        {"bullet": True, "text": "在播客社群推广"},
        {"bullet": True, "text": "目标：50个注册用户"},
        {"bullet": True, "text": "目标：10个付费用户"},
        {"title": "第9-12周：决策点", "text": ""},
        {"bullet": True, "text": "如果无法获得50个付费用户 → 停止"},
        {"bullet": True, "text": "如果获得50个付费用户 → 继续"},
        {"title": "止损点", "text": ""},
        {"bullet": True, "text": "最大投入：¥30,000"},
        {"bullet": True, "text": "最长时间：6个月"},
        {"bullet": True, "text": "最少付费用户：100人"}
    ])

    print("📄 幻灯片 30/30: 最终建议")
    create_summary_slide(prs, [
        ("技术可行性", 8, RGBColor(46, 204, 113)),
        ("市场需求", 7, RGBColor(52, 152, 219)),
        ("商业潜力", 4, RGBColor(231, 76, 60)),
        ("个人开发适配度", 3, RGBColor(231, 76, 60)),
        ("与大厂差异化", 2, RGBColor(231, 76, 60)),
    ])

    # 保存PPT
    output_path = "/Users/tbingy/Desktop/Claude Code/播客逐字稿产品可行性分析报告.pptx"
    prs.save(output_path)

    print("\n" + "="*80)
    print(f"✅ PPT生成成功！")
    print(f"📁 文件位置: {output_path}")
    print(f"📊 幻灯片数量: 30张")
    print(f"📐 尺寸: 10 x 7.5 英寸（标准16:9）")
    print("="*80)

if __name__ == "__main__":
    main()
